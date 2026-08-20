#!/usr/bin/env python3
# ⑤学級開き「大切にしたいこと」見本（8/24研修・触る時間Aの手本）
# 外﨑が実際に使う形。文言はそのまま使っても、自分の言葉に差し替えてもよい。
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

INK = RGBColor(0x15, 0x18, 0x1C)
SUB = RGBColor(0x76, 0x7B, 0x83)
BLUE = RGBColor(0x2B, 0x5F, 0xD9)
NAVY = RGBColor(0x1C, 0x2A, 0x4A)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF2, 0xF2, 0xEF)
LIGHT = RGBColor(0x9F, 0xB0, 0xD0)
FONT = "BIZ UDPゴシック"

W, H = Inches(13.333), Inches(7.5)
OUT = os.path.join(os.path.dirname(__file__), "..", "templates")
os.makedirs(OUT, exist_ok=True)


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(p, bg=PAPER):
    s = p.slides.add_slide(p.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, x, y, w, h, text, size, color=INK, bold=True, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        f = run.font
        f.name = FONT
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
    return tb


def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def photo_slot(s, x, y, w, h, note):
    r = rect(s, x, y, w, h, WASH, line=SUB)
    tf = r.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(note.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(18)
        run.font.bold = False
        run.font.color.rgb = SUB
    return r


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


p = deck()

# 1 表紙
s = blank(p, NAVY)
box(s, 1.0, 2.4, 11.3, 1.5, "2学期 はじめます", 54, PAPER)
box(s, 1.0, 4.1, 11.3, 0.8, "【○年○組】", 26, LIGHT, bold=False)
notes(s, "表紙。学級の名前に差し替え。ここは10秒で次へ。")

# 2 おかえりなさい（夏の写真）
s = blank(p, NAVY)
photo_slot(s, 0.9, 0.7, 7.6, 6.1, "ここに 夏の写真を 1まい\n（学校・行事・空 なんでも\nAIで作った絵でもOK・16:9）")
box(s, 8.9, 2.8, 3.6, 2.2, "おかえり\nなさい", 40, PAPER)
notes(s, "夏の1枚で空気をあたためる。子どもの夏の話を2〜3人分ひろってから、次へ。")

# 3 溜めの1枚
s = blank(p, PAPER)
box(s, 1.0, 2.9, 11.3, 1.6, "先生が 2学期\nいちばん 大切に したいこと", 40, INK,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "ここで一呼吸。「なんだと思う？」と2〜3人に予想させてから次へ進むと、全員がこちらを向きます。")

# 4 メッセージ本体
s = blank(p, PAPER)
photo_slot(s, 0.9, 1.1, 5.3, 5.3, "ここに 絵を 1まい\n（AIで作った 1:1 の絵）\nれい：安心して手をあげる教室")
box(s, 6.8, 2.4, 6.0, 2.6, "「まちがえた」が\n言える教室", 44, INK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "メッセージは1つだけ。自分の学級の言葉に差し替えて。しゃべる中身はスライドに書かず、この1枚を出したまま自分の言葉で話します。")

# 5 なぜか
s = blank(p, PAPER)
rect(s, 1.0, 2.55, 0.18, 1.9, BLUE)
box(s, 1.5, 2.6, 10.8, 1.9, "まちがいは みんなの べんきょうに なるから", 34, INK,
    anchor=MSO_ANCHOR.MIDDLE)
notes(s, "理由も1行だけ。実際のエピソード（1学期にあった場面）をここで1つ話すと、言葉が届きます。")

# 6 問いかけ
s = blank(p, PAPER)
box(s, 1.0, 2.5, 11.3, 1.6, "あなたが 2学期 大切に したいことは？", 40, INK,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.0, 4.5, 11.3, 0.7, "ペアで 30びょう → めあてカードへ", 20, SUB, bold=False,
    align=PP_ALIGN.CENTER)
notes(s, "先生の話で終わらせず、子どもに返す1枚。ここからめあてカードの記入につなぎます。")

# 7 締め
s = blank(p, NAVY)
box(s, 1.0, 2.9, 11.3, 1.5, "2学期も よろしく", 50, PAPER,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "締め。全部で7枚・話す時間は10分が目安です。")

p.save(os.path.join(OUT, "t5_gakkyu_biraki.pptx"))
print("done: t5_gakkyu_biraki.pptx")
