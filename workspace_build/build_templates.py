#!/usr/bin/env python3
# 立岩小研修「文言を差し替えるだけ」スライドテンプレ4本を生成する
# 方針: 16:9 / UDフォント指定 / 色は3色以内 / 差し替え箇所は【 】つき・灰色
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

INK = RGBColor(0x15, 0x18, 0x1C)
SUB = RGBColor(0x76, 0x7B, 0x83)
BLUE = RGBColor(0x2B, 0x5F, 0xD9)
NAVY = RGBColor(0x1C, 0x2A, 0x4A)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF2, 0xF2, 0xEF)
MOON = RGBColor(0xF2, 0xE9, 0xC9)
LIGHT = RGBColor(0x9F, 0xB0, 0xD0)  # navy上の補足文字
FONT = "BIZ UDPゴシック"

W, H = Inches(13.333), Inches(7.5)
OUT = os.path.join(os.path.dirname(__file__), "..", "templates")
os.makedirs(OUT, exist_ok=True)


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(p, bg=PAPER):
    s = p.slides.add_slide(p.slide_layouts[6])  # blank
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, x, y, w, h, text, size, color=INK, bold=True, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        f = run.font
        f.name = FONT
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
    return tb


def rect(s, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def photo_slot(s, x, y, w, h, note="ここに 写真を 1まい\n（この四角を消して、画像をドラッグ）"):
    r = rect(s, x, y, w, h, WASH, line=SUB)
    tf = r.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(note.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(18)
        run.font.bold = False
        run.font.color.rgb = SUB
    return r


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---------- ① 写真ドン型（単元の導入・6枚） ----------
p = deck()

s = blank(p, NAVY)
box(s, 1.0, 2.5, 11.3, 1.6, "【単元名を ここに】", 54, PAPER, align=PP_ALIGN.LEFT)
box(s, 1.0, 4.2, 11.3, 0.8, "【教科 ・ 学年】", 24, LIGHT, bold=False)
notes(s, "表紙。単元名と学年だけ。ここで長く話さず、次の写真へ10秒で進むのがコツです。")

s = blank(p, NAVY)
photo_slot(s, 0.9, 0.7, 7.6, 6.1)
box(s, 8.9, 2.6, 3.6, 2.5, "気づいたこと、\nある？", 40, PAPER)
notes(s, "この1枚が導入の主役。しゃべらず10秒待つ→ペアで30秒→全体で拾う。写真は自分で撮った1枚がいちばん食いつきます。")

s = blank(p, PAPER)
box(s, 1.0, 2.8, 11.3, 2.0, "【問いを ひとつだけ ここに】", 44, INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.0, 5.2, 11.3, 0.7, "よそうを ノートに かこう", 20, SUB, bold=False, align=PP_ALIGN.CENTER)
notes(s, "問いは1つだけ。例:「今夜の月は、どこに見える？」。2つ目の問いが浮かんだら、それは次の時間の分です。")

s = blank(p, PAPER)
rect(s, 1.0, 2.55, 0.18, 1.9, BLUE)
box(s, 1.5, 2.6, 10.8, 1.9, "めあて：【子どもの言葉が出てから ここに】", 36, INK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "めあては子どもの「あれ？」が言葉になってから提示（または黒板に手書き）。スライドを先に見せると、考える前に答え合わせになります。")

s = blank(p, PAPER)
box(s, 1.0, 0.8, 11.3, 1.0, "きょうの ながれ", 28, SUB)
for i, t in enumerate(["【やること 1】", "【やること 2】", "【やること 3】"]):
    y = 2.0 + i * 1.5
    rect(s, 1.0, y, 11.3, 1.15, WASH)
    box(s, 1.35, y + 0.12, 0.9, 0.9, str(i + 1), 32, BLUE)
    box(s, 2.4, y + 0.14, 9.6, 0.9, t, 28, INK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "見通しは3つまで。文言は動詞で(「しらべる」「くらべる」「まとめる」)。")

s = blank(p, PAPER)
box(s, 1.0, 2.8, 11.3, 1.4, "きょう いちばん「あれ？」と思ったことは？", 36, INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.0, 4.6, 11.3, 0.7, "ノートに 2〜3行", 20, SUB, bold=False, align=PP_ALIGN.CENTER)
notes(s, "ふりかえりの観点も1つだけ。導入の時間はここで閉じます。")

p.save(os.path.join(OUT, "t1_shashin_don.pptx"))

# ---------- ② 1枚1ルール型（学級ルール・7枚） ----------
p = deck()

s = blank(p, PAPER)
box(s, 1.0, 2.4, 11.3, 1.5, "2学期、はじめます。", 54, INK)
box(s, 1.0, 4.1, 11.3, 0.8, "【○年○組】", 26, SUB, bold=False)
notes(s, "表紙。始業式の日の学活の1枚目。元気な声でどうぞ。")

s = blank(p, PAPER)
photo_slot(s, 0.9, 0.7, 7.6, 6.1, "ここに 夏の写真を 1まい\n（学校・行事・空 なんでも）")
box(s, 8.9, 2.8, 3.6, 2.2, "おかえり\nなさい。", 40, INK)
notes(s, "夏の1枚で空気をあたためる。子どもの夏の話を2〜3人分ひろってから、次へ。")

for i, (rule, why) in enumerate([
    ("【ルール 1 を ここに】", "れい：チャイムの前に、すわる。"),
    ("【ルール 2 を ここに】", "れい：話している人の ほうを 見る。"),
    ("【ルール 3 を ここに】", "れい：タブレットは 先生の あいずで。"),
]):
    s = blank(p, PAPER)
    box(s, 0.7, 1.6, 3.2, 4.2, str(i + 1), 170, BLUE, align=PP_ALIGN.CENTER)
    rect(s, 3.9, 1.9, 0.02, 3.6, RGBColor(0xE6, 0xE6, 0xE3))
    box(s, 4.4, 2.5, 8.2, 2.0, rule, 40, INK, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 4.4, 4.7, 8.2, 0.7, why, 18, SUB, bold=False)
    box(s, 4.4, 5.6, 8.2, 0.6, f"2学期も つづける やくそく（{i + 1}／3）", 16, SUB, bold=False)
    notes(s, "1枚に1ルール。出したら読み上げず、子どもに読ませる→「なんで あるんだっけ？」を1人に聞く。")

s = blank(p, PAPER)
box(s, 1.0, 2.7, 11.3, 1.6, "なんで、この3つ なんだっけ？", 42, INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.0, 4.6, 11.3, 0.7, "ペアで 30びょう", 20, SUB, bold=False, align=PP_ALIGN.CENTER)
notes(s, "ルールの意味を子どもの言葉で言わせる1枚。ここが「提示」と「押しつけ」の分かれ目です。")

s = blank(p, PAPER)
rect(s, 1.0, 1.0, 0.18, 1.2, BLUE)
box(s, 1.5, 1.05, 10.8, 1.2, "2学期の めあてを かこう", 34, INK, anchor=MSO_ANCHOR.MIDDLE)
for i, t in enumerate(["がんばりたい こと（1つ）", "そのために まいにち すること", "できたか たしかめる 日（月末）"]):
    y = 2.9 + i * 1.3
    rect(s, 1.0, y, 11.3, 1.0, WASH)
    box(s, 1.35, y + 0.1, 0.8, 0.8, str(i + 1), 26, BLUE)
    box(s, 2.3, y + 0.12, 9.7, 0.8, t, 24, INK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "めあてカードへの橋渡し。3つの観点を残したまま、子どもは自分の言葉で書きます。")

p.save(os.path.join(OUT, "t2_gakkyu_rule.pptx"))

# ---------- ③ 授業の流れ・指示型（5枚） ----------
p = deck()

s = blank(p, PAPER)
rect(s, 1.0, 2.75, 0.18, 1.9, BLUE)
box(s, 1.5, 2.8, 10.8, 1.9, "めあて：【ここに】", 40, INK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "授業開始の1枚。黒板にめあてを書く学級は、この1枚を飛ばしてOK。")

s = blank(p, PAPER)
box(s, 1.0, 0.8, 11.3, 1.0, "やること", 28, SUB)
for i, t in enumerate(["【ひとりで かんがえる】", "【ペアで はなす】", "【ノートに まとめる】"]):
    y = 2.0 + i * 1.5
    rect(s, 1.0, y, 11.3, 1.15, WASH)
    box(s, 1.35, y + 0.12, 0.9, 0.9, str(i + 1), 32, BLUE)
    box(s, 2.4, y + 0.14, 9.6, 0.9, t, 28, INK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "活動の前に出す1枚。口頭の指示とちがって、途中で忘れた子が自分で確認できます。出しっぱなしにするのがコツ。")

s = blank(p, PAPER)
box(s, 1.6, 2.0, 4.8, 3.4, "10分", 96, BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rect(s, 6.65, 2.3, 0.02, 2.8, RGBColor(0xE6, 0xE6, 0xE3))
box(s, 7.2, 2.6, 5.0, 2.4, "ゴール：\n【ノートに 3行】", 34, INK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "時間とゴールをセットで。「あと何分？」「どこまでやるの？」の質問がゼロになります。時間は差し替えて。")

s = blank(p, PAPER)
box(s, 1.0, 0.8, 11.3, 1.0, "ペア・グループの うごき", 28, SUB)
for i, t in enumerate(["つくえを あわせる", "きく人は、あいての ほうを 見る", "おわったら、手を あげる"]):
    y = 2.0 + i * 1.5
    rect(s, 1.0, y, 11.3, 1.15, WASH)
    box(s, 1.35, y + 0.12, 0.9, 0.9, str(i + 1), 32, BLUE)
    box(s, 2.4, y + 0.14, 9.6, 0.9, t, 28, INK, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "グループ活動の約束は毎回同じ文言で。同じスライドを使い回すほど、指示が短くなります。")

s = blank(p, PAPER)
box(s, 1.0, 2.4, 11.3, 1.4, "かたづけ → ふりかえり", 40, INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.0, 4.2, 11.3, 1.2, "きょうの「なるほど」を 1つ、ノートに。", 24, SUB, bold=False, align=PP_ALIGN.CENTER)
notes(s, "終末の定型。この1枚が出たら片づけ、が習慣になると授業の終わりが崩れません。")

p.save(os.path.join(OUT, "t3_jugyo_nagare.pptx"))

# ---------- ④ 3択クイズ型（8枚・2問分） ----------
p = deck()

s = blank(p, NAVY)
box(s, 1.0, 2.7, 11.3, 1.5, "クイズで はじめます。", 50, PAPER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.0, 4.5, 11.3, 0.7, "【単元名 ・ 教科】", 22, LIGHT, bold=False, align=PP_ALIGN.CENTER)
notes(s, "表紙。導入でも、朝の会の1問でも。")


def quiz_q(p, qnum, qtext):
    s = blank(p, PAPER)
    box(s, 1.0, 0.7, 2.0, 1.0, f"Q{qnum}", 40, BLUE)
    box(s, 1.0, 1.8, 11.3, 1.6, qtext, 36, INK)
    photo_slot(s, 3.4, 3.6, 6.5, 3.3, "ここに 写真（なくてもOK）")
    notes(s, "問いを読んだら、すぐ次の3択へ。")
    s = blank(p, PAPER)
    labels = ["A", "B", "C"]
    for i, t in enumerate(["【せんたくし A】", "【せんたくし B】", "【せんたくし C】"]):
        x = 0.9 + i * 4.0
        rect(s, x, 2.2, 3.6, 3.0, RGBColor(0xEE, 0xF2, 0xFD), line=BLUE)
        box(s, x, 2.6, 3.6, 1.2, labels[i], 48, BLUE, align=PP_ALIGN.CENTER)
        box(s, x + 0.2, 3.9, 3.2, 1.1, t, 22, INK, align=PP_ALIGN.CENTER)
    notes(s, "指で選ばせる・立たせる・タブレットで送らせる、どれでも。全員が1回「自分の予想」を持つのが目的です。")
    s = blank(p, PAPER)
    box(s, 1.0, 2.2, 11.3, 1.5, "こたえ：【ここに】", 48, BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 1.0, 4.0, 11.3, 1.4, "わけ：【ひとことで ここに】", 26, INK, bold=False, align=PP_ALIGN.CENTER)
    notes(s, "答えより「わけ」が本体。「なんでだと思う？」を先に子どもに聞いてから出すと、導入として強くなります。")


quiz_q(p, 1, "【問題文を ここに】")
quiz_q(p, 2, "【問題文を ここに】")

s = blank(p, NAVY)
box(s, 1.0, 2.7, 11.3, 1.5, "つづきは、この単元で。", 44, PAPER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "クイズの「もっと知りたい」を単元のめあてに接続して閉じます。")

p.save(os.path.join(OUT, "t4_quiz_donyu.pptx"))

print("done:", sorted(os.listdir(OUT)))
